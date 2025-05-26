#!/usr/bin/env python3
"""
Content Verification Test Script

This script verifies that all content from slide_details_updated.json 
has been properly applied to the generated PowerPoint presentation.

Usage: python3 test_content_verification.py
"""

import json
import sys
import os
from pptx import Presentation
from typing import List, Dict, Tuple

# Configuration
UPDATED_JSON_PATH = "./content/slide_details_updated.json"
OUTPUT_PPTX_PATH = "./output/template_updated.pptx"

def load_expected_content(json_path: str) -> List[Dict]:
    """Load the expected content from the JSON file."""
    try:
        with open(json_path, 'r') as f:
            return json.load(f)
    except Exception as e:
        print(f"Error loading JSON file '{json_path}': {e}")
        return []

def extract_text_from_slide(slide) -> List[str]:
    """Extract all text content from a slide."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            # Clean and normalize the text
            text = shape.text.strip()
            if text:
                texts.append(text)
    return texts

def normalize_text(text: str) -> str:
    """Normalize text for comparison by removing extra whitespace and converting to lowercase."""
    if not text:
        return ""
    # Replace multiple whitespaces/newlines with single spaces
    import re
    normalized = re.sub(r'\s+', ' ', text.strip().lower())
    return normalized

def text_similarity_check(expected: str, actual_texts: List[str], threshold: float = 0.8) -> Tuple[bool, str]:
    """
    Check if expected text is found in actual texts with similarity matching.
    Returns (found, best_match)
    """
    if not expected or not actual_texts:
        return False, ""
    
    expected_norm = normalize_text(expected)
    
    # First try exact match
    for actual in actual_texts:
        actual_norm = normalize_text(actual)
        if expected_norm == actual_norm:
            return True, actual
    
    # Then try substring match
    for actual in actual_texts:
        actual_norm = normalize_text(actual)
        if expected_norm in actual_norm or actual_norm in expected_norm:
            return True, actual
    
    # Finally try word-based similarity
    expected_words = set(expected_norm.split())
    best_match = ""
    best_score = 0
    
    for actual in actual_texts:
        actual_norm = normalize_text(actual)
        actual_words = set(actual_norm.split())
        
        if not expected_words or not actual_words:
            continue
            
        # Calculate Jaccard similarity
        intersection = len(expected_words.intersection(actual_words))
        union = len(expected_words.union(actual_words))
        similarity = intersection / union if union > 0 else 0
        
        if similarity > best_score:
            best_score = similarity
            best_match = actual
    
    return best_score >= threshold, best_match

def verify_content_application(json_path: str, pptx_path: str) -> Dict:
    """
    Verify that all content from JSON has been applied to the PowerPoint file.
    Returns a dictionary with verification results.
    """
    results = {
        "total_slides": 0,
        "total_text_shapes": 0,
        "successfully_applied": 0,
        "missing_content": [],
        "slides_with_issues": [],
        "success_rate": 0.0
    }
    
    # Load expected content
    expected_slides = load_expected_content(json_path)
    if not expected_slides:
        print("❌ Failed to load expected content from JSON")
        return results
    
    # Load PowerPoint presentation
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"❌ Failed to load PowerPoint file '{pptx_path}': {e}")
        return results
    
    results["total_slides"] = len(expected_slides)
    
    print(f"🔍 Verifying content application...")
    print(f"📊 Expected slides: {len(expected_slides)}")
    print(f"📊 Actual slides: {len(prs.slides)}")
    print()
    
    # Verify each slide
    for slide_data in expected_slides:
        slide_index = slide_data.get("slide_index", 0)
        slide_layout = slide_data.get("slide_layout_name", "Unknown")
        
        if slide_index >= len(prs.slides):
            print(f"❌ Slide {slide_index + 1} ({slide_layout}): Slide not found in presentation")
            results["slides_with_issues"].append(slide_index + 1)
            continue
        
        actual_slide = prs.slides[slide_index]
        actual_texts = extract_text_from_slide(actual_slide)
        
        slide_missing_content = []
        slide_text_shapes = 0
        slide_applied_content = 0
        
        # Check each shape with text content
        for shape_data in slide_data.get("shapes", []):
            expected_text = shape_data.get("text")
            if not expected_text or not expected_text.strip():
                continue
            
            placeholder_type = shape_data.get("placeholder_type", "Unknown")
            
            # Skip slide number placeholders - they are not applied by design
            if placeholder_type == "SLIDE_NUMBER":
                continue
                
            slide_text_shapes += 1
            results["total_text_shapes"] += 1
            
            shape_name = shape_data.get("name", "Unknown")
            
            # Check if this content was applied
            found, best_match = text_similarity_check(expected_text, actual_texts)
            
            if found:
                slide_applied_content += 1
                results["successfully_applied"] += 1
            else:
                missing_info = {
                    "slide_number": slide_index + 1,
                    "slide_layout": slide_layout,
                    "placeholder_type": placeholder_type,
                    "shape_name": shape_name,
                    "expected_text": expected_text[:100] + "..." if len(expected_text) > 100 else expected_text,
                    "available_texts": [text[:50] + "..." if len(text) > 50 else text for text in actual_texts]
                }
                slide_missing_content.append(missing_info)
                results["missing_content"].append(missing_info)
        
        # Report slide status
        if slide_missing_content:
            results["slides_with_issues"].append(slide_index + 1)
            print(f"⚠️  Slide {slide_index + 1} ({slide_layout}): {slide_applied_content}/{slide_text_shapes} content applied")
            for missing in slide_missing_content:
                print(f"   ❌ Missing {missing['placeholder_type']}: {missing['expected_text']}")
        else:
            if slide_text_shapes > 0:
                print(f"✅ Slide {slide_index + 1} ({slide_layout}): All {slide_text_shapes} content applied")
            else:
                print(f"ℹ️  Slide {slide_index + 1} ({slide_layout}): No text content to verify")
    
    # Calculate success rate
    if results["total_text_shapes"] > 0:
        results["success_rate"] = (results["successfully_applied"] / results["total_text_shapes"]) * 100
    
    return results

def print_summary_report(results: Dict):
    """Print a comprehensive summary report."""
    print("\n" + "="*60)
    print("📋 CONTENT VERIFICATION SUMMARY REPORT")
    print("="*60)
    
    print(f"📊 Total slides verified: {results['total_slides']}")
    print(f"📊 Total text shapes: {results['total_text_shapes']}")
    print(f"✅ Successfully applied: {results['successfully_applied']}")
    print(f"❌ Missing content: {len(results['missing_content'])}")
    print(f"📈 Success rate: {results['success_rate']:.1f}%")
    
    if results["slides_with_issues"]:
        print(f"\n⚠️  Slides with issues: {', '.join(map(str, results['slides_with_issues']))}")
    
    if results["missing_content"]:
        print(f"\n❌ DETAILED MISSING CONTENT:")
        print("-" * 40)
        for i, missing in enumerate(results["missing_content"], 1):
            print(f"{i}. Slide {missing['slide_number']} ({missing['slide_layout']})")
            print(f"   Type: {missing['placeholder_type']}")
            print(f"   Shape: {missing['shape_name']}")
            print(f"   Expected: {missing['expected_text']}")
            if missing['available_texts']:
                print(f"   Available: {missing['available_texts']}")
            else:
                print(f"   Available: [No text found on slide]")
            print()
    
    # Overall status
    if results["success_rate"] == 100:
        print("🎉 RESULT: ALL CONTENT SUCCESSFULLY APPLIED!")
    elif results["success_rate"] >= 90:
        print("✅ RESULT: MOSTLY SUCCESSFUL (Minor issues)")
    elif results["success_rate"] >= 70:
        print("⚠️  RESULT: PARTIALLY SUCCESSFUL (Some issues)")
    else:
        print("❌ RESULT: SIGNIFICANT ISSUES DETECTED")

def main():
    """Main function to run the content verification test."""
    print("🧪 Content Verification Test Script")
    print("=" * 50)
    
    # Check if files exist
    if not os.path.exists(UPDATED_JSON_PATH):
        print(f"❌ JSON file not found: {UPDATED_JSON_PATH}")
        print("Please run main.py first to generate the content.")
        sys.exit(1)
    
    if not os.path.exists(OUTPUT_PPTX_PATH):
        print(f"❌ PowerPoint file not found: {OUTPUT_PPTX_PATH}")
        print("Please run main.py first to generate the presentation.")
        sys.exit(1)
    
    print(f"📁 JSON file: {UPDATED_JSON_PATH}")
    print(f"📁 PPTX file: {OUTPUT_PPTX_PATH}")
    print()
    
    # Run verification
    results = verify_content_application(UPDATED_JSON_PATH, OUTPUT_PPTX_PATH)
    
    # Print summary
    print_summary_report(results)
    
    # Exit with appropriate code
    if results["success_rate"] == 100:
        sys.exit(0)  # Success
    else:
        sys.exit(1)  # Issues found

if __name__ == "__main__":
    main() 