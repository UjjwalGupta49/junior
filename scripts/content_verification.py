import json
import os
from typing import Dict, List, Tuple, Optional
from dataclasses import dataclass
from pptx import Presentation
from extract_slide_details import extract_slide_details
from advanced_placeholder_matcher import AdvancedPlaceholderMatcher, apply_advanced_content_matching

@dataclass
class ContentMismatch:
    """Represents a content verification issue"""
    slide_index: int
    shape_name: str
    placeholder_type: str
    issue_type: str  # 'missing_content', 'default_text', 'content_mismatch', 'empty_placeholder'
    intended_content: str
    actual_content: str
    severity: str  # 'critical', 'warning', 'info'
    description: str

@dataclass
class RepairResult:
    """Results of content repair attempt"""
    mismatch: ContentMismatch
    repair_attempted: bool
    repair_successful: bool
    repair_method: str
    new_content: str
    error_message: Optional[str] = None

@dataclass
class VerificationResult:
    """Results of content verification"""
    total_slides: int
    total_shapes_checked: int
    successful_matches: int
    mismatches: List[ContentMismatch]
    success_rate: float
    overall_status: str  # 'pass', 'warning', 'fail'
    repair_results: Optional[List[RepairResult]] = None
    post_repair_mismatches: Optional[List[ContentMismatch]] = None
    post_repair_success_rate: Optional[float] = None

class ContentVerifier:
    """
    Verifies that content was correctly applied to presentations by comparing
    the final presentation content with the intended content from JSON files.
    """
    
    def __init__(self, debug: bool = False):
        self.debug = debug
        self.default_text_patterns = [
            "click to edit",
            "click to add",
            "add your text here",
            "type your text here",
            "enter your content here",
            "placeholder text",
            "sample text",
            "lorem ipsum",
            "content placeholder",
            "add content",
            "your content here",
            "edit this text",
            "replace this text",
            "template text"
        ]
    
    def verify_presentation_content(self, 
                                  presentation_path: str, 
                                  intended_json_path: str,
                                  auto_fix_critical: bool = False) -> VerificationResult:
        """
        Main verification method that compares presentation content with intended content.
        
        Args:
            presentation_path: Path to the final presentation file
            intended_json_path: Path to the JSON file with intended content
            auto_fix_critical: Whether to attempt automatic fixing of critical issues
            
        Returns:
            VerificationResult object with detailed analysis
        """
        print(f"\n🔍 CONTENT VERIFICATION")
        print(f"   Presentation: {os.path.basename(presentation_path)}")
        print(f"   Source JSON: {os.path.basename(intended_json_path)}")
        if auto_fix_critical:
            print(f"   Auto-fix: ENABLED (will attempt to fix critical issues)")
        print("-" * 50)
        
        try:
            # Extract actual content from presentation
            actual_content = self._extract_actual_content(presentation_path)
            if not actual_content:
                return self._create_failed_result("Failed to extract content from presentation")
            
            # Load intended content from JSON
            intended_content = self._load_intended_content(intended_json_path)
            if not intended_content:
                return self._create_failed_result("Failed to load intended content from JSON")
            
            # Perform detailed comparison
            mismatches = self._compare_content(actual_content, intended_content)
            
            # Calculate success metrics
            total_shapes = sum(len(slide.get("shapes", [])) for slide in intended_content)
            content_shapes = self._count_content_shapes(intended_content)
            successful_matches = content_shapes - len(mismatches)
            success_rate = (successful_matches / content_shapes * 100) if content_shapes > 0 else 0
            
            # Determine overall status
            overall_status = self._determine_overall_status(mismatches, success_rate)
            
            result = VerificationResult(
                total_slides=len(intended_content),
                total_shapes_checked=content_shapes,
                successful_matches=successful_matches,
                mismatches=mismatches,
                success_rate=success_rate,
                overall_status=overall_status
            )
            
            # Auto-fix critical issues if requested
            if auto_fix_critical and mismatches:
                critical_issues = [m for m in mismatches if m.severity == "critical"]
                if critical_issues:
                    print(f"\n🔧 AUTO-REPAIR MODE")
                    print(f"   Found {len(critical_issues)} critical issues to fix")
                    print("-" * 50)
                    
                    repair_results = self._fix_critical_issues(
                        presentation_path, intended_content, critical_issues
                    )
                    result.repair_results = repair_results
                    
                    # Re-verify after repairs
                    post_repair_result = self._post_repair_verification(
                        presentation_path, intended_json_path
                    )
                    result.post_repair_mismatches = post_repair_result.mismatches
                    result.post_repair_success_rate = post_repair_result.success_rate
            
            self._print_verification_summary(result)
            return result
            
        except Exception as e:
            print(f"❌ Error during content verification: {e}")
            return self._create_failed_result(f"Verification error: {e}")
    
    def _extract_actual_content(self, presentation_path: str) -> Optional[List[Dict]]:
        """Extract content from the actual presentation file"""
        try:
            if self.debug:
                print(f"   📄 Extracting content from: {presentation_path}")
            
            actual_content = extract_slide_details(presentation_path)
            
            if actual_content:
                if self.debug:
                    print(f"   ✅ Extracted content from {len(actual_content)} slides")
                return actual_content
            else:
                print(f"   ❌ Failed to extract content from presentation")
                return None
                
        except Exception as e:
            print(f"   ❌ Error extracting actual content: {e}")
            return None
    
    def _load_intended_content(self, json_path: str) -> Optional[List[Dict]]:
        """Load intended content from JSON file"""
        try:
            if self.debug:
                print(f"   📋 Loading intended content from: {json_path}")
                
            with open(json_path, 'r', encoding='utf-8') as f:
                intended_content = json.load(f)
            
            if self.debug:
                print(f"   ✅ Loaded intended content for {len(intended_content)} slides")
            
            return intended_content
            
        except Exception as e:
            print(f"   ❌ Error loading intended content: {e}")
            return None
    
    def _compare_content(self, actual_content: List[Dict], intended_content: List[Dict]) -> List[ContentMismatch]:
        """Compare actual vs intended content and identify mismatches"""
        mismatches = []
        
        if self.debug:
            print(f"   🔄 Comparing {len(actual_content)} actual slides with {len(intended_content)} intended slides")
        
        # Create mapping of slides by index for comparison
        actual_slides = {slide.get("slide_index", i): slide for i, slide in enumerate(actual_content)}
        intended_slides = {slide.get("slide_index", i): slide for i, slide in enumerate(intended_content)}
        
        # Compare each intended slide with actual slide
        for slide_index, intended_slide in intended_slides.items():
            actual_slide = actual_slides.get(slide_index)
            
            if not actual_slide:
                # Missing slide
                for shape in intended_slide.get("shapes", []):
                    if self._is_content_shape(shape):
                        mismatches.append(ContentMismatch(
                            slide_index=slide_index,
                            shape_name=shape.get("name", "Unknown"),
                            placeholder_type=shape.get("placeholder_type", "Unknown"),
                            issue_type="missing_slide",
                            intended_content=shape.get("text", ""),
                            actual_content="",
                            severity="critical",
                            description=f"Entire slide {slide_index} is missing"
                        ))
                continue
            
            # Compare shapes within the slide
            slide_mismatches = self._compare_slide_shapes(slide_index, actual_slide, intended_slide)
            mismatches.extend(slide_mismatches)
        
        return mismatches
    
    def _compare_slide_shapes(self, slide_index: int, actual_slide: Dict, intended_slide: Dict) -> List[ContentMismatch]:
        """Compare shapes between actual and intended slide"""
        mismatches = []
        
        actual_shapes = {shape.get("name", ""): shape for shape in actual_slide.get("shapes", [])}
        intended_shapes = intended_slide.get("shapes", [])
        
        for intended_shape in intended_shapes:
            if not self._is_content_shape(intended_shape):
                continue  # Skip non-content shapes
            
            shape_name = intended_shape.get("name", "")
            placeholder_type = intended_shape.get("placeholder_type", "")
            intended_text = intended_shape.get("text", "").strip()
            
            # Skip empty intended content
            if not intended_text or intended_text in ["", "(missing)"]:
                continue
            
            # Find matching actual shape
            actual_shape = self._find_matching_actual_shape(shape_name, placeholder_type, actual_shapes, intended_shape)
            
            if not actual_shape:
                mismatches.append(ContentMismatch(
                    slide_index=slide_index,
                    shape_name=shape_name,
                    placeholder_type=placeholder_type,
                    issue_type="missing_shape",
                    intended_content=intended_text,
                    actual_content="",
                    severity="critical",
                    description=f"Shape '{shape_name}' not found in actual presentation"
                ))
                continue
            
            actual_text = actual_shape.get("text", "").strip()
            
            # Check for various content issues
            mismatch = self._analyze_content_match(slide_index, shape_name, placeholder_type, intended_text, actual_text)
            if mismatch:
                mismatches.append(mismatch)
        
        return mismatches
    
    def _find_matching_actual_shape(self, shape_name: str, placeholder_type: str, actual_shapes: Dict, intended_shape: Dict) -> Optional[Dict]:
        """Find the best matching actual shape for an intended shape"""
        # Try exact name match first
        if shape_name in actual_shapes:
            return actual_shapes[shape_name]
        
        # Try matching by placeholder type and position
        intended_left = intended_shape.get("left_inches", 0)
        intended_top = intended_shape.get("top_inches", 0)
        
        best_match = None
        best_score = float('inf')
        
        for actual_shape in actual_shapes.values():
            actual_type = actual_shape.get("placeholder_type", "")
            actual_left = actual_shape.get("left_inches", 0)
            actual_top = actual_shape.get("top_inches", 0)
            
            # Calculate match score (lower is better)
            score = 0
            
            # Type match bonus
            if actual_type == placeholder_type:
                score -= 100
            
            # Position distance penalty
            if intended_left is not None and actual_left is not None:
                score += abs(intended_left - actual_left) * 10
            if intended_top is not None and actual_top is not None:
                score += abs(intended_top - actual_top) * 10
            
            if score < best_score:
                best_score = score
                best_match = actual_shape
        
        return best_match if best_score < 50 else None  # Reasonable threshold
    
    def _analyze_content_match(self, slide_index: int, shape_name: str, placeholder_type: str, intended_text: str, actual_text: str) -> Optional[ContentMismatch]:
        """Analyze if content matches and identify specific issues"""
        
        # Normalize text for comparison (remove extra whitespace, convert to lowercase)
        intended_normalized = " ".join(intended_text.split()).lower()
        actual_normalized = " ".join(actual_text.split()).lower()
        
        # Check for empty actual content
        if not actual_text:
            return ContentMismatch(
                slide_index=slide_index,
                shape_name=shape_name,
                placeholder_type=placeholder_type,
                issue_type="empty_placeholder",
                intended_content=intended_text,
                actual_content=actual_text,
                severity="critical",
                description="Placeholder is empty - no content was applied"
            )
        
        # Check for default template text
        if self._is_default_text(actual_text):
            return ContentMismatch(
                slide_index=slide_index,
                shape_name=shape_name,
                placeholder_type=placeholder_type,
                issue_type="default_text",
                intended_content=intended_text,
                actual_content=actual_text,
                severity="critical",
                description=f"Contains default template text: '{actual_text[:50]}...'"
            )
        
        # Check for exact content match
        if intended_normalized == actual_normalized:
            return None  # Perfect match, no mismatch
        
        # Check for partial match (intended content is subset of actual)
        if intended_normalized in actual_normalized or actual_normalized in intended_normalized:
            similarity = self._calculate_similarity(intended_normalized, actual_normalized)
            if similarity > 0.8:  # 80% similarity threshold
                return None  # Close enough match
            else:
                return ContentMismatch(
                    slide_index=slide_index,
                    shape_name=shape_name,
                    placeholder_type=placeholder_type,
                    issue_type="content_mismatch",
                    intended_content=intended_text,
                    actual_content=actual_text,
                    severity="warning",
                    description=f"Content partially matches (similarity: {similarity:.1%})"
                )
        
        # Complete content mismatch
        return ContentMismatch(
            slide_index=slide_index,
            shape_name=shape_name,
            placeholder_type=placeholder_type,
            issue_type="content_mismatch",
            intended_content=intended_text,
            actual_content=actual_text,
            severity="critical",
            description="Content completely different from intended"
        )
    
    def _is_content_shape(self, shape: Dict) -> bool:
        """Check if a shape should have content applied (exclude slide numbers, etc.)"""
        placeholder_type = shape.get("placeholder_type", "")
        return placeholder_type not in ["SLIDE_NUMBER"] and shape.get("has_text_frame", False)
    
    def _is_default_text(self, text: str) -> bool:
        """Check if text appears to be default template text"""
        text_lower = text.lower().strip()
        
        # Check against known default text patterns
        for pattern in self.default_text_patterns:
            if pattern in text_lower:
                return True
        
        # Check for very generic/short text that might be default
        if len(text_lower) < 10 and any(word in text_lower for word in ["click", "add", "edit", "text", "content"]):
            return True
        
        return False
    
    def _calculate_similarity(self, text1: str, text2: str) -> float:
        """Calculate text similarity between two strings"""
        if not text1 or not text2:
            return 0.0
        
        # Simple word-based similarity
        words1 = set(text1.split())
        words2 = set(text2.split())
        
        if not words1 or not words2:
            return 0.0
        
        intersection = len(words1.intersection(words2))
        union = len(words1.union(words2))
        
        return intersection / union if union > 0 else 0.0
    
    def _count_content_shapes(self, intended_content: List[Dict]) -> int:
        """Count total number of content shapes that should be verified"""
        count = 0
        for slide in intended_content:
            for shape in slide.get("shapes", []):
                if self._is_content_shape(shape) and shape.get("text", "").strip():
                    count += 1
        return count
    
    def _determine_overall_status(self, mismatches: List[ContentMismatch], success_rate: float) -> str:
        """Determine overall verification status"""
        critical_issues = sum(1 for m in mismatches if m.severity == "critical")
        warning_issues = sum(1 for m in mismatches if m.severity == "warning")
        
        if critical_issues == 0 and warning_issues == 0:
            return "pass"
        elif critical_issues == 0 and warning_issues <= 2:
            return "warning"
        else:
            return "fail"
    
    def _print_verification_summary(self, result: VerificationResult):
        """Print a detailed verification summary"""
        print(f"\n📊 VERIFICATION RESULTS")
        print(f"   Total slides: {result.total_slides}")
        print(f"   Content shapes checked: {result.total_shapes_checked}")
        print(f"   Successful matches: {result.successful_matches}")
        print(f"   Success rate: {result.success_rate:.1f}%")
        print(f"   Overall status: {result.overall_status.upper()}")
        
        if result.mismatches:
            print(f"\n⚠️  ISSUES FOUND ({len(result.mismatches)}):")
            
            # Group by severity
            critical = [m for m in result.mismatches if m.severity == "critical"]
            warnings = [m for m in result.mismatches if m.severity == "warning"]
            
            if critical:
                print(f"   🚨 Critical Issues ({len(critical)}):")
                for mismatch in critical[:5]:  # Show first 5
                    print(f"      • Slide {mismatch.slide_index + 1}, {mismatch.shape_name}: {mismatch.description}")
                if len(critical) > 5:
                    print(f"      ... and {len(critical) - 5} more critical issues")
            
            if warnings:
                print(f"   ⚠️  Warnings ({len(warnings)}):")
                for mismatch in warnings[:3]:  # Show first 3
                    print(f"      • Slide {mismatch.slide_index + 1}, {mismatch.shape_name}: {mismatch.description}")
                if len(warnings) > 3:
                    print(f"      ... and {len(warnings) - 3} more warnings")
        else:
            print(f"   ✅ No issues found - all content applied successfully!")
        
        # Show repair results if auto-fix was performed
        if result.repair_results is not None:
            self._print_repair_summary(result.repair_results)
        
        # Show post-repair verification results if available
        if result.post_repair_mismatches is not None and result.post_repair_success_rate is not None:
            self._print_post_repair_summary(result.post_repair_mismatches, result.post_repair_success_rate, result.success_rate)
    
    def _print_repair_summary(self, repair_results: List[RepairResult]):
        """Print detailed repair attempt summary"""
        print(f"\n🔧 REPAIR ATTEMPT SUMMARY")
        print("-" * 30)
        
        total_repairs = len(repair_results)
        attempted_repairs = sum(1 for r in repair_results if r.repair_attempted)
        successful_repairs = sum(1 for r in repair_results if r.repair_successful)
        
        print(f"   Total critical issues: {total_repairs}")
        print(f"   Repair attempts: {attempted_repairs}")
        print(f"   Successful repairs: {successful_repairs}")
        
        if attempted_repairs > 0:
            repair_success_rate = (successful_repairs / attempted_repairs) * 100
            print(f"   Repair success rate: {repair_success_rate:.1f}%")
        
        # Group repairs by method/strategy
        repair_methods = {}
        for repair in repair_results:
            if repair.repair_successful:
                method = repair.repair_method
                if method not in repair_methods:
                    repair_methods[method] = 0
                repair_methods[method] += 1
        
        if repair_methods:
            print(f"\n   🎯 Successful repair strategies:")
            for method, count in sorted(repair_methods.items(), key=lambda x: x[1], reverse=True):
                print(f"      • {method}: {count} fixes")
        
        # Show failed repairs
        failed_repairs = [r for r in repair_results if r.repair_attempted and not r.repair_successful]
        if failed_repairs:
            print(f"\n   ❌ Failed repairs ({len(failed_repairs)}):")
            for repair in failed_repairs[:3]:  # Show first 3
                shape_name = repair.mismatch.shape_name
                slide_num = repair.mismatch.slide_index + 1
                error = repair.error_message or "Unknown error"
                print(f"      • Slide {slide_num}, {shape_name}: {error}")
            if len(failed_repairs) > 3:
                print(f"      ... and {len(failed_repairs) - 3} more failed repairs")
    
    def _print_post_repair_summary(self, 
                                 post_repair_mismatches: List[ContentMismatch], 
                                 post_repair_success_rate: float,
                                 original_success_rate: float):
        """Print post-repair verification summary"""
        print(f"\n📈 POST-REPAIR VERIFICATION")
        print("-" * 30)
        
        improvement = post_repair_success_rate - original_success_rate
        remaining_critical = sum(1 for m in post_repair_mismatches if m.severity == "critical")
        remaining_warnings = sum(1 for m in post_repair_mismatches if m.severity == "warning")
        
        print(f"   Original success rate: {original_success_rate:.1f}%")
        print(f"   Post-repair success rate: {post_repair_success_rate:.1f}%")
        print(f"   Improvement: {improvement:+.1f}%")
        
        if improvement > 0:
            print(f"   ✅ Repair process improved content verification!")
        elif improvement == 0:
            print(f"   ➖ No improvement from repair process")
        else:
            print(f"   ❌ Repair process may have caused issues")
        
        print(f"\n   Remaining issues: {len(post_repair_mismatches)}")
        if remaining_critical > 0:
            print(f"      🚨 Critical: {remaining_critical}")
        if remaining_warnings > 0:
            print(f"      ⚠️  Warnings: {remaining_warnings}")
        
        if len(post_repair_mismatches) == 0:
            print(f"   🎉 All critical issues successfully resolved!")
        elif remaining_critical == 0:
            print(f"   ✅ All critical issues resolved (only warnings remain)")
    
    def _create_failed_result(self, error_message: str) -> VerificationResult:
        """Create a failed verification result"""
        return VerificationResult(
            total_slides=0,
            total_shapes_checked=0,
            successful_matches=0,
            mismatches=[],
            success_rate=0.0,
            overall_status="fail"
        )
    
    def _fix_critical_issues(self, 
                           presentation_path: str, 
                           intended_content: List[Dict], 
                           critical_issues: List[ContentMismatch]) -> List[RepairResult]:
        """
        Attempt to fix critical issues by applying correct content to problematic shapes.
        
        Args:
            presentation_path: Path to the presentation file
            intended_content: JSON data with intended slide data
            critical_issues: List of critical issues to fix
            
        Returns:
            List of RepairResult objects showing what was attempted and results
        """
        repair_results = []
        
        try:
            # Load presentation for editing
            prs = Presentation(presentation_path)
            
            # Group issues by slide for efficient processing
            issues_by_slide = {}
            for issue in critical_issues:
                slide_idx = issue.slide_index
                if slide_idx not in issues_by_slide:
                    issues_by_slide[slide_idx] = []
                issues_by_slide[slide_idx].append(issue)
            
            print(f"   📝 Processing {len(issues_by_slide)} slides with critical issues")
            
            # Process each slide with issues
            for slide_idx, slide_issues in issues_by_slide.items():
                if slide_idx < len(prs.slides) and slide_idx < len(intended_content):
                    slide = prs.slides[slide_idx]
                    slide_data = intended_content[slide_idx]
                    
                    print(f"\n   🔧 Repairing slide {slide_idx + 1} ({len(slide_issues)} issues)")
                    
                    # Attempt repairs for this slide
                    slide_repairs = self._repair_slide_issues(
                        slide, slide_data, slide_issues
                    )
                    repair_results.extend(slide_repairs)
            
            # Save the repaired presentation
            prs.save(presentation_path)
            print(f"\n   💾 Saved repaired presentation: {os.path.basename(presentation_path)}")
            
        except Exception as e:
            print(f"   ❌ Error during repair process: {e}")
            # Add failed repair results for all issues
            for issue in critical_issues:
                repair_results.append(RepairResult(
                    mismatch=issue,
                    repair_attempted=False,
                    repair_successful=False,
                    repair_method="none",
                    new_content="",
                    error_message=str(e)
                ))
        
        return repair_results
    
    def _repair_slide_issues(self, 
                           slide, 
                           slide_data: Dict, 
                           slide_issues: List[ContentMismatch]) -> List[RepairResult]:
        """
        Repair issues for a specific slide using advanced placeholder matching.
        
        Args:
            slide: PowerPoint slide object
            slide_data: JSON data for the slide
            slide_issues: List of issues for this slide
            
        Returns:
            List of RepairResult objects
        """
        repair_results = []
        
        # Create a matcher for this slide
        matcher = AdvancedPlaceholderMatcher(debug=self.debug)
        
        # Get all intended shapes for this slide
        intended_shapes = slide_data.get("shapes", [])
        
        # Create a mapping of shape names to intended content
        intended_content_map = {}
        for shape in intended_shapes:
            shape_name = shape.get("name", "")
            shape_text = shape.get("text", "")
            if shape_name and shape_text and shape_text.strip():
                intended_content_map[shape_name] = shape
        
        # Attempt to repair each issue
        for issue in slide_issues:
            repair_result = self._attempt_single_repair(
                slide, slide_data, issue, intended_content_map, matcher
            )
            repair_results.append(repair_result)
        
        return repair_results
    
    def _attempt_single_repair(self, 
                             slide, 
                             slide_data: Dict, 
                             issue: ContentMismatch, 
                             intended_content_map: Dict,
                             matcher: AdvancedPlaceholderMatcher) -> RepairResult:
        """
        Attempt to repair a single content issue.
        
        Args:
            slide: PowerPoint slide object
            slide_data: JSON data for the slide
            issue: The specific issue to repair
            intended_content_map: Mapping of shape names to intended content
            matcher: Advanced placeholder matcher instance
            
        Returns:
            RepairResult object
        """
        shape_name = issue.shape_name
        intended_content = issue.intended_content
        
        try:
            # Find the intended shape data
            intended_shape = intended_content_map.get(shape_name)
            if not intended_shape:
                return RepairResult(
                    mismatch=issue,
                    repair_attempted=False,
                    repair_successful=False,
                    repair_method="no_intended_content",
                    new_content="",
                    error_message=f"No intended content found for shape '{shape_name}'"
                )
            
            # Use the advanced matcher to find and apply content
            matches = matcher.match_placeholders_for_slide(slide, {"shapes": [intended_shape]})
            
            if not matches:
                return RepairResult(
                    mismatch=issue,
                    repair_attempted=True,
                    repair_successful=False,
                    repair_method="no_match_found",
                    new_content="",
                    error_message=f"Could not find matching shape for '{shape_name}'"
                )
            
            # Apply the content using the best match
            best_match = matches[0]
            ppt_shape = best_match.ppt_shape
            new_text = intended_shape.get("text", "").strip()
            
            if not new_text:
                return RepairResult(
                    mismatch=issue,
                    repair_attempted=True,
                    repair_successful=False,
                    repair_method="empty_content",
                    new_content="",
                    error_message="Intended content is empty"
                )
            
            # Apply the content
            ppt_shape.text_frame.clear()
            
            # Handle multi-line text
            lines = new_text.split('\n')
            for i, line in enumerate(lines):
                line = line.strip()
                if not line:
                    continue
                
                if i == 0:
                    if ppt_shape.text_frame.paragraphs:
                        ppt_shape.text_frame.paragraphs[0].text = line
                    else:
                        p = ppt_shape.text_frame.add_paragraph()
                        p.text = line
                else:
                    p = ppt_shape.text_frame.add_paragraph()
                    p.text = line
            
            print(f"      ✅ Fixed '{shape_name}' using {best_match.strategy.value}")
            print(f"         Applied: '{new_text[:50]}{'...' if len(new_text) > 50 else ''}'")
            
            return RepairResult(
                mismatch=issue,
                repair_attempted=True,
                repair_successful=True,
                repair_method=best_match.strategy.value,
                new_content=new_text,
                error_message=None
            )
            
        except Exception as e:
            print(f"      ❌ Failed to fix '{shape_name}': {e}")
            return RepairResult(
                mismatch=issue,
                repair_attempted=True,
                repair_successful=False,
                repair_method="error",
                new_content="",
                error_message=str(e)
            )
    
    def _post_repair_verification(self, 
                                presentation_path: str, 
                                intended_json_path: str) -> VerificationResult:
        """
        Perform verification after repairs to check improvement.
        
        Args:
            presentation_path: Path to the repaired presentation
            intended_json_path: Path to the JSON file with intended content
            
        Returns:
            VerificationResult with post-repair status
        """
        print(f"\n   🔍 POST-REPAIR VERIFICATION")
        print("-" * 30)
        
        # Create a new verifier instance without auto-fix to avoid recursion
        verifier = ContentVerifier(debug=False)
        result = verifier.verify_presentation_content(
            presentation_path, intended_json_path, auto_fix_critical=False
        )
        
        return result

def verify_presentation_content(presentation_path: str, 
                               intended_json_path: str, 
                               debug: bool = False,
                               auto_fix_critical: bool = False) -> VerificationResult:
    """
    Convenience function to verify presentation content.
    
    Args:
        presentation_path: Path to the final presentation file
        intended_json_path: Path to the JSON file with intended content
        debug: Enable debug output
        auto_fix_critical: Whether to attempt automatic fixing of critical issues
        
    Returns:
        VerificationResult object
    """
    verifier = ContentVerifier(debug=debug)
    return verifier.verify_presentation_content(
        presentation_path, intended_json_path, auto_fix_critical=auto_fix_critical
    )

if __name__ == "__main__":
    # Test the verification module
    import sys
    import argparse
    
    parser = argparse.ArgumentParser(
        description="Content Verification with Auto-Fix",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python content_verification.py presentation.pptx content.json                    # Basic verification
  python content_verification.py presentation.pptx content.json --debug           # With debug output
  python content_verification.py presentation.pptx content.json --auto-fix        # With auto-fix
  python content_verification.py presentation.pptx content.json --auto-fix --debug # Full debug + auto-fix
        """
    )
    
    parser.add_argument("presentation_path", help="Path to the presentation file")
    parser.add_argument("intended_json_path", help="Path to the JSON file with intended content")
    parser.add_argument("--debug", action="store_true", help="Enable debug output")
    parser.add_argument("--auto-fix", action="store_true", help="Automatically fix critical issues")
    
    args = parser.parse_args()
    
    print("🔍 Testing Content Verification with Auto-Fix")
    print("=" * 60)
    
    if args.auto_fix:
        print("🔧 Auto-fix mode ENABLED - will attempt to repair critical issues")
    else:
        print("📊 Verification-only mode - no repairs will be attempted")
    
    print("=" * 60)
    
    result = verify_presentation_content(
        args.presentation_path, 
        args.intended_json_path, 
        debug=args.debug,
        auto_fix_critical=args.auto_fix
    )
    
    print(f"\n🎯 Final Status: {result.overall_status.upper()}")
    if result.overall_status == "pass":
        print("✅ Content verification passed!")
    elif result.overall_status == "warning":
        print("⚠️  Content verification passed with warnings")
    else:
        print("❌ Content verification failed")
    
    # Show repair summary if auto-fix was used
    if args.auto_fix and result.repair_results:
        successful_repairs = sum(1 for r in result.repair_results if r.repair_successful)
        total_critical = len(result.repair_results)
        print(f"\n🔧 Auto-fix completed: {successful_repairs}/{total_critical} critical issues resolved")
        
        if result.post_repair_success_rate is not None:
            improvement = result.post_repair_success_rate - result.success_rate
            print(f"📈 Overall improvement: {improvement:+.1f}% success rate") 