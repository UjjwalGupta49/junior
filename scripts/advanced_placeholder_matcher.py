import json
import re
from typing import Dict, List, Tuple, Optional, Set
from pptx import Presentation
from pptx.shapes.base import BaseShape
from dataclasses import dataclass
from enum import Enum

class MatchingStrategy(Enum):
    EXACT_NAME = "exact_name"
    PLACEHOLDER_NUMBER = "placeholder_number"
    POSITION_BASED = "position_based"
    LAYOUT_PATTERN = "layout_pattern"
    CONTENT_SIMILARITY = "content_similarity"
    FALLBACK_ORDER = "fallback_order"

@dataclass
class PlaceholderMatch:
    """Represents a match between JSON shape data and PowerPoint shape"""
    json_shape: Dict
    ppt_shape: BaseShape
    confidence: float
    strategy: MatchingStrategy
    reasoning: str

@dataclass
class SlideContext:
    """Context information about a slide to help with matching"""
    slide_index: int
    layout_name: str
    total_content_placeholders: int
    content_placeholder_positions: List[Tuple[float, float]]  # (left, top) positions
    available_shapes: List[BaseShape]
    used_shapes: Set[int]

class AdvancedPlaceholderMatcher:
    """
    Advanced matcher for handling multiple content placeholders on the same slide.
    Uses multiple strategies and contextual information for accurate matching.
    """
    
    def __init__(self, debug: bool = False):
        self.debug = debug
        self.position_tolerance = 500000  # EMU units (about 0.5 inch)
        
    def match_placeholders_for_slide(self, slide, slide_data: Dict) -> List[PlaceholderMatch]:
        """
        Main method to match all placeholders for a single slide.
        Returns a list of PlaceholderMatch objects sorted by confidence.
        """
        # Create slide context
        context = self._create_slide_context(slide, slide_data)
        
        # Get all shapes that need content (TITLE and OBJECT only)
        content_shapes = self._filter_content_shapes(slide_data)
        
        if self.debug:
            print(f"\n🔍 Matching slide {context.slide_index + 1} ({context.layout_name})")
            print(f"   Content shapes to match: {len(content_shapes)}")
            print(f"   Available PPT shapes: {len(context.available_shapes)}")
        
        matches = []
        
        # Apply matching strategies in order of preference
        for shape_data in content_shapes:
            best_match = self._find_best_match(shape_data, context)
            if best_match:
                matches.append(best_match)
                context.used_shapes.add(id(best_match.ppt_shape))
                
                if self.debug:
                    print(f"   ✅ Matched '{shape_data.get('name')}' -> '{best_match.ppt_shape.name}' "
                          f"(confidence: {best_match.confidence:.2f}, strategy: {best_match.strategy.value})")
            else:
                if self.debug:
                    print(f"   ❌ No match found for '{shape_data.get('name')}'")
        
        return sorted(matches, key=lambda m: m.confidence, reverse=True)
    
    def _create_slide_context(self, slide, slide_data: Dict) -> SlideContext:
        """Create context information about the slide"""
        layout_name = slide_data.get("slide_layout_name", "Unknown")
        
        # Count content placeholders in JSON
        content_count = sum(1 for shape in slide_data.get("shapes", []) 
                          if shape.get("placeholder_type") == "OBJECT")
        
        # Get positions of content placeholders
        positions = []
        for shape in slide_data.get("shapes", []):
            if shape.get("placeholder_type") == "OBJECT":
                left = shape.get("left_inches", 0)
                top = shape.get("top_inches", 0)
                positions.append((left, top))
        
        # Get available shapes from PowerPoint slide
        available_shapes = [shape for shape in slide.shapes 
                          if shape.has_text_frame and self._is_content_shape(shape)]
        
        return SlideContext(
            slide_index=slide_data.get("slide_index", 0),
            layout_name=layout_name,
            total_content_placeholders=content_count,
            content_placeholder_positions=positions,
            available_shapes=available_shapes,
            used_shapes=set()
        )
    
    def _filter_content_shapes(self, slide_data: Dict) -> List[Dict]:
        """Filter shapes to exclude only SLIDE_NUMBER placeholders - include all others"""
        content_shapes = []
        for shape_data in slide_data.get("shapes", []):
            placeholder_type = shape_data.get("placeholder_type", "")
            
            # Skip only SLIDE_NUMBER placeholders - include everything else
            if placeholder_type != "SLIDE_NUMBER":
                content_shapes.append(shape_data)
        
        # Sort by priority: TITLE first, then OBJECT, then others
        def priority(shape):
            placeholder_type = shape.get("placeholder_type", "")
            if placeholder_type == "TITLE":
                return 0
            elif placeholder_type == "OBJECT":
                return 1
            else:
                return 2  # All other types (SUBTITLE, BODY, etc.)
        
        return sorted(content_shapes, key=priority)
    
    def _is_content_shape(self, shape: BaseShape) -> bool:
        """Check if a PowerPoint shape is a content shape (not slide number, picture, etc.)"""
        if not shape.has_text_frame:
            return False
        
        name_lower = shape.name.lower()
        text = shape.text.strip()
        
        # Exclude obvious non-content shapes
        exclusions = [
            "picture" in name_lower,
            "table" in name_lower,
            ("slide" in name_lower and "number" in name_lower),
            (text.isdigit() and len(text) <= 2),  # Likely slide numbers
        ]
        
        return not any(exclusions)
    
    def _find_best_match(self, shape_data: Dict, context: SlideContext) -> Optional[PlaceholderMatch]:
        """Find the best matching PowerPoint shape for a JSON shape"""
        placeholder_type = shape_data.get("placeholder_type", "")
        shape_name = shape_data.get("name", "")
        
        # Try different matching strategies
        strategies = [
            self._try_exact_name_match,
            self._try_placeholder_number_match,
            self._try_position_match,
            self._try_layout_pattern_match,
            self._try_content_similarity_match,
            self._try_fallback_order_match
        ]
        
        best_match = None
        best_confidence = 0.0
        
        for strategy_func in strategies:
            match = strategy_func(shape_data, context)
            if match and match.confidence > best_confidence:
                best_match = match
                best_confidence = match.confidence
                
                # If we get a very high confidence match, use it immediately
                if best_confidence >= 0.9:
                    break
        
        return best_match
    
    def _try_exact_name_match(self, shape_data: Dict, context: SlideContext) -> Optional[PlaceholderMatch]:
        """Strategy 1: Exact name matching"""
        shape_name = shape_data.get("name", "")
        if not shape_name:
            return None
        
        for shape in context.available_shapes:
            if (id(shape) not in context.used_shapes and 
                shape.name and 
                shape.name.lower() == shape_name.lower()):
                return PlaceholderMatch(
                    json_shape=shape_data,
                    ppt_shape=shape,
                    confidence=1.0,
                    strategy=MatchingStrategy.EXACT_NAME,
                    reasoning=f"Exact name match: '{shape_name}'"
                )
        return None
    
    def _try_placeholder_number_match(self, shape_data: Dict, context: SlideContext) -> Optional[PlaceholderMatch]:
        """Strategy 2: Match by placeholder number (e.g., 'Content Placeholder 1')"""
        shape_name = shape_data.get("name", "")
        placeholder_type = shape_data.get("placeholder_type", "")
        
        if placeholder_type != "OBJECT" or "placeholder" not in shape_name.lower():
            return None
        
        # Extract number from JSON shape name
        json_match = re.search(r'placeholder\s+(\d+)', shape_name.lower())
        if not json_match:
            return None
        
        target_number = int(json_match.group(1))
        
        # Find PowerPoint shape with matching number
        for shape in context.available_shapes:
            if id(shape) not in context.used_shapes:
                ppt_match = re.search(r'placeholder\s+(\d+)', shape.name.lower())
                if ppt_match and int(ppt_match.group(1)) == target_number:
                    return PlaceholderMatch(
                        json_shape=shape_data,
                        ppt_shape=shape,
                        confidence=0.95,
                        strategy=MatchingStrategy.PLACEHOLDER_NUMBER,
                        reasoning=f"Placeholder number match: {target_number}"
                    )
        return None
    
    def _try_position_match(self, shape_data: Dict, context: SlideContext) -> Optional[PlaceholderMatch]:
        """Strategy 3: Position-based matching with tight tolerance"""
        left_pos = shape_data.get("left_inches", 0)
        top_pos = shape_data.get("top_inches", 0)
        placeholder_type = shape_data.get("placeholder_type", "")
        
        best_match = None
        best_distance = float('inf')
        
        for shape in context.available_shapes:
            if id(shape) not in context.used_shapes:
                try:
                    distance = abs(shape.left - left_pos) + abs(shape.top - top_pos)
                    
                    if distance < best_distance and distance < self.position_tolerance:
                        # Additional validation for OBJECT placeholders
                        if placeholder_type == "OBJECT":
                            if "content" in shape.name.lower() and "placeholder" in shape.name.lower():
                                best_distance = distance
                                best_match = shape
                        elif placeholder_type == "TITLE":
                            if "title" in shape.name.lower():
                                best_distance = distance
                                best_match = shape
                        else:
                            best_distance = distance
                            best_match = shape
                except:
                    continue
        
        if best_match:
            confidence = max(0.1, 1.0 - (best_distance / self.position_tolerance))
            return PlaceholderMatch(
                json_shape=shape_data,
                ppt_shape=best_match,
                confidence=confidence,
                strategy=MatchingStrategy.POSITION_BASED,
                reasoning=f"Position match (distance: {best_distance:.0f})"
            )
        return None
    
    def _try_layout_pattern_match(self, shape_data: Dict, context: SlideContext) -> Optional[PlaceholderMatch]:
        """Strategy 4: Match based on layout patterns and slide structure"""
        placeholder_type = shape_data.get("placeholder_type", "")
        
        if placeholder_type == "TITLE":
            # Find title shapes
            for shape in context.available_shapes:
                if (id(shape) not in context.used_shapes and 
                    "title" in shape.name.lower()):
                    return PlaceholderMatch(
                        json_shape=shape_data,
                        ppt_shape=shape,
                        confidence=0.8,
                        strategy=MatchingStrategy.LAYOUT_PATTERN,
                        reasoning="Title placeholder pattern match"
                    )
        
        elif placeholder_type == "SUBTITLE":
            # Find subtitle shapes
            for shape in context.available_shapes:
                if (id(shape) not in context.used_shapes and 
                    "subtitle" in shape.name.lower()):
                    return PlaceholderMatch(
                        json_shape=shape_data,
                        ppt_shape=shape,
                        confidence=0.8,
                        strategy=MatchingStrategy.LAYOUT_PATTERN,
                        reasoning="Subtitle placeholder pattern match"
                    )
        
        elif placeholder_type == "OBJECT":
            # For content placeholders, use layout-specific logic
            layout_name = context.layout_name.lower()
            
            # Get available content placeholders
            content_shapes = [shape for shape in context.available_shapes 
                            if (id(shape) not in context.used_shapes and 
                                "content" in shape.name.lower() and 
                                "placeholder" in shape.name.lower())]
            
            if content_shapes:
                # Sort by position (left to right, top to bottom)
                content_shapes.sort(key=lambda s: (s.top, s.left))
                
                # For layouts with multiple content areas, use positional logic
                if "two content" in layout_name or "comparison" in layout_name:
                    shape_name = shape_data.get("name", "").lower()
                    
                    # Try to match based on left/right positioning
                    if "1" in shape_name or "left" in shape_name:
                        # Use leftmost shape
                        return PlaceholderMatch(
                            json_shape=shape_data,
                            ppt_shape=content_shapes[0],
                            confidence=0.85,
                            strategy=MatchingStrategy.LAYOUT_PATTERN,
                            reasoning="Left content area in two-content layout"
                        )
                    elif "2" in shape_name or "right" in shape_name:
                        # Use rightmost shape
                        if len(content_shapes) > 1:
                            return PlaceholderMatch(
                                json_shape=shape_data,
                                ppt_shape=content_shapes[-1],
                                confidence=0.85,
                                strategy=MatchingStrategy.LAYOUT_PATTERN,
                                reasoning="Right content area in two-content layout"
                            )
                
                # Default: use first available content shape
                return PlaceholderMatch(
                    json_shape=shape_data,
                    ppt_shape=content_shapes[0],
                    confidence=0.7,
                    strategy=MatchingStrategy.LAYOUT_PATTERN,
                    reasoning="First available content placeholder"
                )
        
        return None
    
    def _try_content_similarity_match(self, shape_data: Dict, context: SlideContext) -> Optional[PlaceholderMatch]:
        """Strategy 5: Match based on existing content similarity"""
        expected_text = shape_data.get("text", "")
        if not expected_text:
            return None
        
        best_match = None
        best_similarity = 0.0
        
        for shape in context.available_shapes:
            if id(shape) not in context.used_shapes:
                existing_text = shape.text.strip()
                if existing_text:
                    similarity = self._calculate_text_similarity(expected_text, existing_text)
                    if similarity > best_similarity and similarity > 0.3:
                        best_similarity = similarity
                        best_match = shape
        
        if best_match:
            return PlaceholderMatch(
                json_shape=shape_data,
                ppt_shape=best_match,
                confidence=best_similarity * 0.6,  # Lower confidence for content-based matching
                strategy=MatchingStrategy.CONTENT_SIMILARITY,
                reasoning=f"Content similarity: {best_similarity:.2f}"
            )
        return None
    
    def _try_fallback_order_match(self, shape_data: Dict, context: SlideContext) -> Optional[PlaceholderMatch]:
        """Strategy 6: Fallback - use order-based matching for any placeholder type"""
        placeholder_type = shape_data.get("placeholder_type", "")
        
        if placeholder_type == "OBJECT":
            # Find any unused content placeholder
            for shape in context.available_shapes:
                if (id(shape) not in context.used_shapes and 
                    "content" in shape.name.lower() and 
                    "placeholder" in shape.name.lower()):
                    return PlaceholderMatch(
                        json_shape=shape_data,
                        ppt_shape=shape,
                        confidence=0.5,
                        strategy=MatchingStrategy.FALLBACK_ORDER,
                        reasoning="Fallback: first available content placeholder"
                    )
        
        elif placeholder_type == "TITLE":
            # Find any unused title shape
            for shape in context.available_shapes:
                if (id(shape) not in context.used_shapes and 
                    "title" in shape.name.lower()):
                    return PlaceholderMatch(
                        json_shape=shape_data,
                        ppt_shape=shape,
                        confidence=0.6,
                        strategy=MatchingStrategy.FALLBACK_ORDER,
                        reasoning="Fallback: first available title shape"
                    )
        
        elif placeholder_type == "SUBTITLE":
            # Find any unused subtitle shape
            for shape in context.available_shapes:
                if (id(shape) not in context.used_shapes and 
                    "subtitle" in shape.name.lower()):
                    return PlaceholderMatch(
                        json_shape=shape_data,
                        ppt_shape=shape,
                        confidence=0.6,
                        strategy=MatchingStrategy.FALLBACK_ORDER,
                        reasoning="Fallback: first available subtitle shape"
                    )
        
        else:
            # Generic fallback for any other placeholder type
            placeholder_name_lower = placeholder_type.lower()
            for shape in context.available_shapes:
                if (id(shape) not in context.used_shapes and 
                    placeholder_name_lower in shape.name.lower()):
                    return PlaceholderMatch(
                        json_shape=shape_data,
                        ppt_shape=shape,
                        confidence=0.4,
                        strategy=MatchingStrategy.FALLBACK_ORDER,
                        reasoning=f"Fallback: first available {placeholder_type} shape"
                    )
            
            # Last resort: any unused text shape
            for shape in context.available_shapes:
                if id(shape) not in context.used_shapes:
                    return PlaceholderMatch(
                        json_shape=shape_data,
                        ppt_shape=shape,
                        confidence=0.3,
                        strategy=MatchingStrategy.FALLBACK_ORDER,
                        reasoning="Last resort: any available text shape"
                    )
        
        return None
    
    def _calculate_text_similarity(self, text1: str, text2: str) -> float:
        """Calculate similarity between two text strings"""
        if not text1 or not text2:
            return 0.0
        
        # Simple word-based similarity
        words1 = set(text1.lower().split())
        words2 = set(text2.lower().split())
        
        if not words1 or not words2:
            return 0.0
        
        intersection = words1.intersection(words2)
        union = words1.union(words2)
        
        return len(intersection) / len(union) if union else 0.0

def apply_advanced_content_matching(slide, slide_data: Dict, debug: bool = False) -> int:
    """
    Apply content to a slide using advanced placeholder matching.
    Returns the number of shapes successfully updated.
    """
    matcher = AdvancedPlaceholderMatcher(debug=debug)
    matches = matcher.match_placeholders_for_slide(slide, slide_data)
    
    updated_count = 0
    
    for match in matches:
        new_text = match.json_shape.get("text", "")
        if not new_text or "(missing)" in new_text:
            # Clean the text
            if new_text:
                new_text = new_text.replace(" (missing)", "").replace("(missing)", "").strip()
            if not new_text:
                continue
        
        try:
            # Apply the content
            shape = match.ppt_shape
            shape.text_frame.clear()
            
            # Handle multi-line text
            lines = new_text.split('\n')
            for i, line in enumerate(lines):
                line = line.strip()
                if not line:
                    continue
                
                if i == 0:
                    if shape.text_frame.paragraphs:
                        shape.text_frame.paragraphs[0].text = line
                    else:
                        p = shape.text_frame.add_paragraph()
                        p.text = line
                else:
                    p = shape.text_frame.add_paragraph()
                    p.text = line
            
            updated_count += 1
            
            if debug:
                print(f"   ✅ Applied content to '{shape.name}': '{new_text[:50]}...'")
                print(f"      Strategy: {match.strategy.value}, Confidence: {match.confidence:.2f}")
                print(f"      Reasoning: {match.reasoning}")
        
        except Exception as e:
            if debug:
                print(f"   ❌ Error applying content to '{match.ppt_shape.name}': {e}")
    
    return updated_count

# Test function
if __name__ == "__main__":
    # Test with sample data
    print("Advanced Placeholder Matcher - Test Mode")
    print("This module provides sophisticated matching for multiple content placeholders.") 