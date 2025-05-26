import json
from pptx import Presentation
from pptx.util import Inches # Though not strictly needed for text update, good to have if positions were also updated

def update_presentation_from_json(json_path, pptx_path):
    """
    Updates text content in a PowerPoint presentation based on a JSON file.
    Preserves existing formatting of shapes.

    Args:
        json_path (str): Path to the JSON file with slide and shape details.
        pptx_path (str): Path to the PowerPoint file to be updated.
    """
    try:
        with open(json_path, 'r') as f:
            all_slides_data = json.load(f)
    except FileNotFoundError:
        print(f"Error: JSON file not found at '{json_path}'")
        return False
    except json.JSONDecodeError:
        print(f"Error: Could not decode JSON from '{json_path}'")
        return False
    except Exception as e:
        print(f"An unexpected error occurred while reading JSON '{json_path}': {e}")
        return False

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error loading presentation '{pptx_path}': {e}")
        return False

    if len(prs.slides) == 0 and len(all_slides_data) > 0:
        print(f"Warning: The presentation '{pptx_path}' is empty, but JSON data exists.")
        # Decide if you want to proceed or return, for now, let's proceed carefully

    updated_slides_count = 0
    updated_shapes_count = 0

    for slide_data in all_slides_data:
        slide_idx = slide_data.get("slide_index")
        if slide_idx is None:
            print(f"Warning: Slide data missing 'slide_index'. Skipping: {slide_data.get('slide_layout_name', 'Unknown slide')}")
            continue

        if slide_idx >= len(prs.slides):
            print(f"Warning: slide_index {slide_idx} from JSON is out of bounds for presentation with {len(prs.slides)} slides. Skipping.")
            continue

        slide = prs.slides[slide_idx]
        shapes_data = slide_data.get("shapes", [])
        slide_updated_this_iteration = False

        for shape_data in shapes_data:
            shape_id = shape_data.get("shape_id")
            new_text = shape_data.get("text") # Expecting new text here

            if shape_id is None:
                print(f"Warning: Shape data on slide {slide_idx} missing 'shape_id'. Skipping shape.")
                continue

            found_shape = None
            for s in slide.shapes:
                if s.shape_id == shape_id:
                    found_shape = s
                    break
            
            if found_shape:
                if found_shape.has_text_frame and new_text is not None: # Only update if new_text is provided
                    # Clear existing content to ensure formatting is based on the shape's style
                    # not carried over from old paragraph runs if they had different styles.
                    found_shape.text_frame.clear() 
                    
                    # Split the input text by newline characters
                    lines = new_text.split('\n')
                    
                    # Handle case where new_text might be an empty string,
                    # which split('\n') results in [''].
                    # If new_text was truly empty, we want one empty paragraph.
                    # If new_text had content, lines will have that content.
                    if not lines and new_text == "": # Handles empty string input correctly
                        p = found_shape.text_frame.add_paragraph()
                        p.text = "" # Explicitly set an empty paragraph
                    else:
                        for line_content in lines:
                            # Each line becomes a new paragraph.
                            # This new paragraph will inherit the style (including bullet formatting)
                            # from the text frame or its placeholder definition in the template.
                            p = found_shape.text_frame.add_paragraph()
                            p.text = line_content
                            
                    updated_shapes_count += 1
                    if not slide_updated_this_iteration:
                        updated_slides_count +=1
                        slide_updated_this_iteration = True
                elif not found_shape.has_text_frame and new_text is not None:
                    print(f"Warning: Shape ID {shape_id} on slide {slide_idx} does not have a text frame. Cannot update text: '{new_text[:50]}...'")
            else:
                print(f"Warning: Shape with ID {shape_id} not found on slide {slide_idx}. Skipping update for this shape.")

    try:
        prs.save(pptx_path)
        print(f"Presentation '{pptx_path}' updated successfully.")
        print(f"Updated text in {updated_shapes_count} shapes across {updated_slides_count} slides.")
        return True
    except Exception as e:
        print(f"Error saving presentation '{pptx_path}': {e}")
        return False

if __name__ == "__main__":
    JSON_INPUT_FILE = "slide_details_updated.json"
    PPTX_TARGET_FILE = "template_updated.pptx"

    print(f"Starting presentation update process...")
    print(f"Reading updates from: {JSON_INPUT_FILE}")
    print(f"Applying updates to: {PPTX_TARGET_FILE}")

    success = update_presentation_from_json(JSON_INPUT_FILE, PPTX_TARGET_FILE)

    if success:
        print("Update process completed successfully.")
    else:
        print("Update process encountered errors.") 