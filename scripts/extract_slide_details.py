from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE # Reverted to simple import
from pptx.enum.dml import MSO_THEME_COLOR # Example, if needed for more detail later
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT # For potential future use
import json

def get_enum_name(enum_value, default_name="UNKNOWN"):
    """Safely gets the name of an enum member, or a default if not found/applicable."""
    if enum_value is None:
        return default_name
    try:
        return enum_value.name
    except AttributeError:
        # Handle cases where enum_value might be an int (e.g. some shape types not in MSO_SHAPE)
        return str(enum_value)

def extract_slide_details(presentation_path):
    """
    Extracts detailed information about each slide in a PowerPoint presentation.
    Includes shapes, their properties, and text content.
    """
    try:
        prs = Presentation(presentation_path)
    except Exception as e:
        print(f"Error loading presentation '{presentation_path}': {e}")
        return None

    all_slides_data = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {
            "slide_index": slide_idx,
            "slide_layout_name": slide.slide_layout.name if slide.slide_layout else "N/A",
            "shapes": []
        }

        for shape in slide.shapes:
            shape_info = {
                "name": shape.name,
                "shape_id": shape.shape_id,
                "shape_type": get_enum_name(shape.shape_type, default_name=str(shape.shape_type)),
                "is_placeholder": shape.is_placeholder,
                "left_inches": round(Inches(shape.left).inches, 3) if shape.left is not None else None,
                "top_inches": round(Inches(shape.top).inches, 3) if shape.top is not None else None,
                "width_inches": round(Inches(shape.width).inches, 3) if shape.width is not None else None,
                "height_inches": round(Inches(shape.height).inches, 3) if shape.height is not None else None,
                "has_text_frame": shape.has_text_frame,
                "text": None,
                "is_table": False,
                "is_chart": False,
                "is_picture": False,
            }

            if shape.is_placeholder:
                try:
                    shape_info["placeholder_type"] = get_enum_name(shape.placeholder_format.type)
                except AttributeError: # placeholder_format might not always be fully populated
                    shape_info["placeholder_type"] = "UNKNOWN_PLACEHOLDER_FORMAT"


            if shape.has_text_frame and shape.text_frame:
                shape_info["text"] = shape.text_frame.text

            # Check for specific shape types using their integer values
            current_shape_type_val = shape.shape_type
            if hasattr(shape.shape_type, 'value'): # If it's an enum object, get its value
                current_shape_type_val = shape.shape_type.value

            # MSO_SHAPE.TABLE.value is 19
            if current_shape_type_val == 19 and shape.has_table:
                shape_info["is_table"] = True
            
            # MSO_SHAPE.CHART.value is 3
            if current_shape_type_val == 3 and shape.has_chart:
                shape_info["is_chart"] = True
            
            # MSO_SHAPE.PICTURE.value is 13
            if current_shape_type_val == 13:
                shape_info["is_picture"] = True
            
            # Could add more specific handlers for other shape types if needed
            # e.g., MSO_SHAPE.AUTO_SHAPE.value is 1

            slide_data["shapes"].append(shape_info)
        
        all_slides_data.append(slide_data)

    return all_slides_data

if __name__ == "__main__":
    # Assuming template.pptx is in the parent directory of this script
    INPUT_PPTX_FILE = "./template.pptx" 
    OUTPUT_JSON_FILE = "./slide_details.json"

    print(f"Extracting slide details from: {INPUT_PPTX_FILE}")
    slide_info = extract_slide_details(INPUT_PPTX_FILE)

    if slide_info:
        try:
            with open(OUTPUT_JSON_FILE, "w") as f:
                json.dump(slide_info, f, indent=4)
            print(f"Successfully extracted slide details to: {OUTPUT_JSON_FILE}")
            print(f"Found details for {len(slide_info)} slides.")
        except IOError as e:
            print(f"Error writing JSON to file '{OUTPUT_JSON_FILE}': {e}")
        except Exception as e: # Catch other potential errors during JSON dump
            print(f"An unexpected error occurred during JSON serialization: {e}")
    else:
        print(f"No slide information was extracted from '{INPUT_PPTX_FILE}'. This could be due to an error loading the file or an empty presentation.") 