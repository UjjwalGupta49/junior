import json
import shutil
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from advanced_placeholder_matcher import apply_advanced_content_matching

def create_and_update_reorganized_presentation(original_pptx_path: str, updated_json_path: str, output_pptx_path: str) -> bool:
    """
    Creates a new presentation from updated JSON data by using slide layouts and applying content directly.
    This approach avoids shape ID mismatches by building slides from scratch.
    
    Args:
        original_pptx_path: Path to the original template presentation
        updated_json_path: Path to the updated slides JSON with new content
        output_pptx_path: Path where the new presentation will be saved
    
    Returns:
        True if successful, False otherwise
    """
    try:
        with open(updated_json_path, 'r') as f:
            slides_data = json.load(f)
    except Exception as e:
        print(f"Error reading updated JSON: {e}")
        return False
    
    try:
        # Load the original presentation to access layouts
        template_prs = Presentation(original_pptx_path)
        
        # Create a new empty presentation using the same template
        new_prs = Presentation(original_pptx_path)
        
        # Remove all existing slides to start fresh
        slide_count = len(new_prs.slides)
        for i in range(slide_count - 1, -1, -1):
            slide_id = new_prs.slides[i].slide_id
            new_prs.part.drop_rel(new_prs.slides._sldIdLst[i].rId)
            del new_prs.slides._sldIdLst[i]
        
        print("--- Creating reorganized presentation structure ---")
        
        # Create slides from JSON data using appropriate layouts
        for slide_data in slides_data:
            layout_name = slide_data.get("slide_layout_name", "Title and Content")
            
            # Find the matching layout
            slide_layout = find_layout_by_name(template_prs, layout_name)
            if not slide_layout:
                print(f"Warning: Layout '{layout_name}' not found, using default")
                slide_layout = template_prs.slide_layouts[1]  # Default to content layout
            
            # Add new slide with the layout
            new_slide = new_prs.slides.add_slide(slide_layout)
            
            # Apply content to the slide using advanced matching
            updated_shapes = apply_advanced_content_matching(new_slide, slide_data, debug=True)
            print(f"Updated {updated_shapes} shapes in slide '{slide_data.get('slide_layout_name', 'Unknown')}'")
            
            # Fallback to original method if advanced matching fails
            if updated_shapes == 0:
                print("  Falling back to original matching method...")
                apply_content_to_slide(new_slide, slide_data)
        
        # Save the new presentation
        new_prs.save(output_pptx_path)
        print(f"Successfully created reorganized presentation with {len(slides_data)} slides")
        
        return True
        
    except Exception as e:
        print(f"Error creating reorganized presentation: {e}")
        return False

def find_layout_by_name(presentation, layout_name):
    """
    Finds a slide layout by name in the presentation.
    """
    try:
        for layout in presentation.slide_layouts:
            if layout.name.lower() == layout_name.lower():
                return layout
        
        # If exact match not found, try partial matching
        for layout in presentation.slide_layouts:
            if layout_name.lower() in layout.name.lower() or layout.name.lower() in layout_name.lower():
                return layout
                
    except Exception as e:
        print(f"Warning: Error finding layout '{layout_name}': {e}")
    
    return None

def apply_content_to_slide(slide, slide_data):
    """
    Applies content from JSON data to a slide by matching placeholders and shapes.
    This approach uses placeholder types and positions rather than shape IDs.
    """
    try:
        shapes_data = slide_data.get("shapes", [])
        updated_shapes = 0
        used_shapes = set()  # Track used shapes by their object id
        
        # Sort shapes by priority: TITLE first, then OBJECT (main content), then SLIDE_NUMBER last
        def shape_priority(shape_data):
            placeholder_type = shape_data.get("placeholder_type", "")
            if placeholder_type == "TITLE":
                return 0  # Highest priority
            elif placeholder_type == "OBJECT":
                return 1  # Second priority (main content)
            elif placeholder_type == "SLIDE_NUMBER":
                return 9  # Lowest priority (should be applied last)
            else:
                return 2  # Other types
        
        sorted_shapes_data = sorted(shapes_data, key=shape_priority)
        
        for shape_data in sorted_shapes_data:
            placeholder_type = shape_data.get("placeholder_type", "")
            shape_type = shape_data.get("shape_type", "")
            shape_name = shape_data.get("name", "")
            
            # Skip slide number placeholders - do not apply content to them
            if placeholder_type == "SLIDE_NUMBER":
                print(f"  Skipping SLIDE_NUMBER placeholder (content not applied)")
                continue
            
            # Apply content to all placeholder types except SLIDE_NUMBER
            # This makes the script work with multiple templates that may have different placeholder types
            
            # Debug: Show what we're processing
            print(f"  Processing {placeholder_type} placeholder: {shape_name} (shape_type: {shape_type})")
            
            new_text = shape_data.get("text")
            if not new_text or "(missing)" in new_text:
                # Clean the text by removing (missing) markers
                if new_text:
                    new_text = new_text.replace(" (missing)", "").replace("(missing)", "").strip()
                if not new_text:
                    continue
            
            # Find the best matching shape on the slide
            target_shape = find_best_matching_shape(slide, shape_data, used_shapes)
            
            if target_shape and target_shape.has_text_frame:
                print(f"    -> Found target shape: {target_shape.name}")
                # Mark this shape as used
                used_shapes.add(id(target_shape))
                
                # Clear existing content and add new content
                target_shape.text_frame.clear()
                
                # Handle multi-line text
                lines = new_text.split('\n')
                for i, line in enumerate(lines):
                    line = line.strip()
                    if not line:
                        continue
                        
                    if i == 0:
                        # First line goes in the first paragraph
                        if target_shape.text_frame.paragraphs:
                            target_shape.text_frame.paragraphs[0].text = line
                        else:
                            p = target_shape.text_frame.add_paragraph()
                            p.text = line
                    else:
                        # Additional lines get new paragraphs
                        p = target_shape.text_frame.add_paragraph()
                        p.text = line
                
                updated_shapes += 1
                print(f"  Applied content to {shape_data.get('placeholder_type', 'unknown')} placeholder ({shape_data.get('name', 'unknown')}): '{new_text[:50]}...'")
            else:
                print(f"  ❌ Could not find matching shape for {shape_data.get('placeholder_type', 'unknown')} placeholder ({shape_data.get('name', 'unknown')})")
        
        if updated_shapes > 0:
            print(f"Updated {updated_shapes} shapes in slide '{slide_data.get('slide_layout_name', 'Unknown')}'")
        else:
            print(f"Warning: No shapes updated in slide '{slide_data.get('slide_layout_name', 'Unknown')}'")
            
    except Exception as e:
        print(f"Warning: Error applying content to slide: {e}")

def find_best_matching_shape(slide, shape_data, used_shapes):
    """
    Finds the best matching shape on a slide based on placeholder type, position, and characteristics.
    Enhanced to handle multiple content placeholders on the same slide.
    """
    try:
        placeholder_type = shape_data.get("placeholder_type")
        is_placeholder = shape_data.get("is_placeholder", False)
        shape_name = shape_data.get("name", "")
        left_pos = shape_data.get("left_inches", 0)
        top_pos = shape_data.get("top_inches", 0)
        
        # For debugging multiple content placeholders
        # print(f"    Looking for shape: {shape_name} ({placeholder_type}) at position ({left_pos}, {top_pos})")
        
        # Strategy 0: For OBJECT placeholders, try exact name match first (handles multiple content areas)
        if placeholder_type == "OBJECT" and shape_name:
            for shape in slide.shapes:
                if (id(shape) not in used_shapes and 
                    shape.has_text_frame and 
                    shape.name and 
                    shape.name.lower() == shape_name.lower()):
                    # print(f"    Found exact name match: {shape.name}")
                    return shape
        
        # Strategy 1: Match by placeholder type (improved logic)
        if is_placeholder and placeholder_type:
            # Try exact placeholder type match first
            for shape in slide.shapes:
                if (id(shape) not in used_shapes and 
                    shape.is_placeholder and 
                    hasattr(shape.placeholder_format, 'type')):
                    try:
                        if shape.placeholder_format.type.name == placeholder_type:
                            # print(f"    Found placeholder type match: {shape.name}")
                            return shape
                    except:
                        # If type comparison fails, continue to other strategies
                        pass
            
            # Strategy 1.1: Match SLIDE_NUMBER placeholders (very specific matching)
            if placeholder_type == "SLIDE_NUMBER":
                # First try exact slide number placeholder names
                for shape in slide.shapes:
                    if (id(shape) not in used_shapes and 
                        shape.is_placeholder and 
                        shape.has_text_frame):
                        shape_name_lower = shape.name.lower()
                        if ("slide" in shape_name_lower and "number" in shape_name_lower):
                            return shape
                
                # Then try shapes that already contain numeric content (likely slide numbers)
                for shape in slide.shapes:
                    if (id(shape) not in used_shapes and 
                        shape.is_placeholder and 
                        shape.has_text_frame):
                        text = shape.text.strip()
                        # Very specific criteria for slide numbers
                        if (text.isdigit() and len(text) <= 2 and 
                            int(text) <= 50):  # Reasonable slide number range
                            return shape
            
            # Strategy 1.2: Match TITLE placeholders
            elif placeholder_type == "TITLE":
                for shape in slide.shapes:
                    if (id(shape) not in used_shapes and 
                        shape.is_placeholder and 
                        shape.has_text_frame):
                        # Match title placeholders
                        if ("title" in shape.name.lower() or
                            hasattr(shape.placeholder_format, 'type') and 
                            shape.placeholder_format.type.name == "TITLE"):
                            return shape
            
            # Strategy 1.3: Match OBJECT/CONTENT placeholders (prioritize main content areas)
            elif placeholder_type == "OBJECT":
                # First try content placeholders specifically
                for shape in slide.shapes:
                    if (id(shape) not in used_shapes and 
                        shape.is_placeholder and 
                        shape.has_text_frame):
                        shape_name_lower = shape.name.lower()
                        if ("content" in shape_name_lower and 
                            "placeholder" in shape_name_lower):
                            return shape
                
                # Then try any placeholder that's clearly not title or slide number
                for shape in slide.shapes:
                    if (id(shape) not in used_shapes and 
                        shape.has_text_frame and
                        shape.is_placeholder):
                        shape_name_lower = shape.name.lower()
                        text = shape.text.strip()
                        
                        # Exclude obvious title and slide number shapes more aggressively
                        is_title = "title" in shape_name_lower
                        is_slide_number = (("slide" in shape_name_lower and "number" in shape_name_lower) or
                                         (text.isdigit() and len(text) <= 2))
                        is_picture = "picture" in shape_name_lower
                        is_table = "table" in shape_name_lower
                        
                        if not (is_title or is_slide_number or is_picture or is_table):
                            return shape
        
        # Strategy 2: Match by shape name pattern (improved)
        if shape_name:
            # Try exact name match first
            for shape in slide.shapes:
                if (id(shape) not in used_shapes and 
                    shape.has_text_frame and 
                    shape.name and 
                    shape.name.lower() == shape_name.lower()):
                    return shape
            
            # Try partial name match
            for shape in slide.shapes:
                if (id(shape) not in used_shapes and 
                    shape.has_text_frame and 
                    shape.name and 
                    (shape_name.lower() in shape.name.lower() or 
                     any(word in shape.name.lower() for word in shape_name.lower().split()))):
                    return shape
        
        # Strategy 3: Match by position (with improved tolerance for multiple content areas)
        position_tolerance = 500000  # EMU units (about 0.5 inch) - very tight tolerance for precise matching
        best_match = None
        best_distance = float('inf')
        
        for shape in slide.shapes:
            if (id(shape) not in used_shapes and 
                shape.has_text_frame and
                shape.is_placeholder):
                try:
                    distance = abs(shape.left - left_pos) + abs(shape.top - top_pos)
                    if distance < best_distance and distance < position_tolerance:
                        # Additional check: ensure it's a content placeholder for OBJECT types
                        if placeholder_type == "OBJECT":
                            shape_name_lower = shape.name.lower()
                            if "content" in shape_name_lower and "placeholder" in shape_name_lower:
                                best_distance = distance
                                best_match = shape
                        else:
                            best_distance = distance
                            best_match = shape
                except:
                    continue
        
        if best_match:
            return best_match
        
        # Strategy 4: Smart fallback for multiple content placeholders
        if placeholder_type == "OBJECT":
            # Get all available content placeholders sorted by position (left to right, top to bottom)
            content_placeholders = []
            for shape in slide.shapes:
                if (id(shape) not in used_shapes and 
                    shape.has_text_frame and 
                    shape.is_placeholder):
                    shape_name_lower = shape.name.lower()
                    text = shape.text.strip()
                    
                    # Skip obvious non-content shapes
                    is_title = "title" in shape_name_lower
                    is_slide_number = (("slide" in shape_name_lower and "number" in shape_name_lower) or
                                     (text.isdigit() and len(text) <= 2))
                    is_picture = "picture" in shape_name_lower
                    is_table = "table" in shape_name_lower
                    
                    if not (is_title or is_slide_number or is_picture or is_table):
                        content_placeholders.append(shape)
            
            if content_placeholders:
                # Sort by position: left to right, then top to bottom
                content_placeholders.sort(key=lambda s: (s.top, s.left))
                
                # Try to match by placeholder number in name
                if shape_name and "placeholder" in shape_name.lower():
                    # Extract number from shape name (e.g., "Content Placeholder 1" -> 1)
                    import re
                    match = re.search(r'placeholder\s+(\d+)', shape_name.lower())
                    if match:
                        target_number = int(match.group(1))
                        # Find placeholder with matching number
                        for shape in content_placeholders:
                            shape_match = re.search(r'placeholder\s+(\d+)', shape.name.lower())
                            if shape_match and int(shape_match.group(1)) == target_number:
                                return shape
                
                # Fallback: return the first available content placeholder
                return content_placeholders[0]
        
        elif placeholder_type == "SLIDE_NUMBER":
            # For slide numbers, only match very specific shapes
            for shape in slide.shapes:
                if (id(shape) not in used_shapes and 
                    shape.has_text_frame and 
                    shape.is_placeholder):
                    shape_name_lower = shape.name.lower()
                    text = shape.text.strip()
                    
                    # Only match if it's clearly a slide number
                    if (("slide" in shape_name_lower and "number" in shape_name_lower) or
                        (text.isdigit() and len(text) <= 2 and int(text) <= 50)):
                        return shape
        
        elif placeholder_type == "TITLE":
            for shape in slide.shapes:
                if (id(shape) not in used_shapes and 
                    shape.has_text_frame and 
                    shape.is_placeholder):
                    if "title" in shape.name.lower():
                        return shape
        
        # Strategy 5: Last resort - any unused text shape
        for shape in slide.shapes:
            if (id(shape) not in used_shapes and 
                shape.has_text_frame):
                return shape
                
    except Exception as e:
        print(f"Warning: Error finding matching shape: {e}")
    
    return None

# Legacy functions for backward compatibility
def create_presentation_from_reorganized_json(original_pptx_path: str, reorganized_json_path: str, output_pptx_path: str) -> bool:
    """Legacy function - redirects to the new implementation"""
    return create_and_update_reorganized_presentation(original_pptx_path, reorganized_json_path, output_pptx_path)

def update_reorganized_presentation_content(reorganized_json_path: str, presentation_path: str) -> bool:
    """Legacy function - no longer needed with the new approach"""
    print("Note: update_reorganized_presentation_content is no longer needed with the new approach")
    return True

if __name__ == "__main__":
    # Test the module
    original_pptx = "./template/template.pptx"
    reorganized_json = "./content/slide_details_reorganized.json"
    output_pptx = "./output/template_reorganized.pptx"
    
    success = create_and_update_reorganized_presentation(original_pptx, reorganized_json, output_pptx)
    print(f"Test completed. Success: {success}") 